#!/usr/bin/env python3
"""Generate Heimvorteil investor pitch deck as PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Brand colors
EMERALD = RGBColor(0x05, 0x96, 0x69)
EMERALD_DARK = RGBColor(0x04, 0x73, 0x57)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1C, 0x1C, 0x1C)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tf

# ========== SLIDE 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, EMERALD_DARK)
add_text_box(slide, 1.5, 1.5, 10, 1.5, "HEIMVORTEIL", 60, True, WHITE, PP_ALIGN.LEFT)
add_text_box(slide, 1.5, 3.2, 10, 1, "Netflix für dein Dorf", 36, False, RGBColor(0x6E, 0xE7, 0xB7), PP_ALIGN.LEFT)
add_text_box(slide, 1.5, 4.5, 10, 1, "Subscription-based local commerce membership\nfor DACH villages", 20, False, RGBColor(0xA7, 0xF3, 0xD0), PP_ALIGN.LEFT)
add_text_box(slide, 1.5, 6.2, 10, 0.5, "heimvorteil.at", 16, False, RGBColor(0xA7, 0xF3, 0xD0), PP_ALIGN.LEFT)

# ========== SLIDE 2: Problem ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1.5, 0.8, 10, 0.8, "Das Problem", 40, True, EMERALD_DARK)
add_bullet_list(slide, 1.5, 2.0, 10, 5, [
    "Lokale Geschäfte verlieren Kunden an Online-Handel und Supermarktketten",
    "Dörfer verlieren wirtschaftliche Vitalität — Ortskerne sterben aus",
    "Bestehende Lösungen (Regionalgeld, Payback, Groupon) sind komplex und skalieren nicht",
    "72% aller KMU-Digitalisierungsprojekte scheitern an mangelnder Akzeptanz",
    "€ Milliarden an lokaler Kaufkraft fließen jedes Jahr aus DACH-Gemeinden ab",
], 20)

# ========== SLIDE 3: Solution ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1.5, 0.8, 10, 0.8, "Die Lösung: Heimvorteil", 40, True, EMERALD_DARK)
add_bullet_list(slide, 1.5, 2.0, 10, 5, [
    "Monatsabo (€9,90–24,90) = sofort 15–25% Rabatt bei ALLEN lokalen Geschäften",
    "Keine Gutscheine, keine Punkte — Karte scannen, sparen",
    "Gemeinde bekommt Echtzeit-Dashboard + Dorffonds (10% jedes Abos)",
    "Betriebe zahlen nichts — tauschen Marge gegen garantierten Fußverkehr",
    "Kein Zahlungslizenz nötig — identifikationsbasiert, nicht zahlungsbasiert",
], 20)

# ========== SLIDE 4: Market ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1.5, 0.8, 10, 0.8, "Der Markt", 40, True, EMERALD_DARK)
add_bullet_list(slide, 1.5, 2.0, 5.5, 4.5, [
    "11.000+ Gemeinden in der DACH-Region",
    "2.000+ Dörfer mit 500–5.000 Einwohnern (Sweet Spot)",
    "Hyperlocal Services: $2,9B → $3,3B (15,4% CAGR)",
    "TAM: 2.000 Dörfer × €30k/Jahr = €60M ARR",
], 20)

# ========== SLIDE 5: Business Model ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1.5, 0.8, 10, 0.8, "Geschäftsmodell", 40, True, EMERALD_DARK)
add_text_box(slide, 1.5, 2.0, 5, 0.5, "3 Abo-Stufen", 22, True, DARK)
add_bullet_list(slide, 1.5, 2.7, 5, 2, [
    "Basis: €9,90/Mo — 15% Rabatt",
    "Plus: €14,90/Mo — 20% Rabatt + NFC-Karte",
    "Dorf: €24,90/Mo — 25% Rabatt + Familie",
], 18)
add_text_box(slide, 1.5, 4.5, 5, 0.5, "Revenue Split", 22, True, DARK)
add_bullet_list(slide, 1.5, 5.2, 5, 2, [
    "70% → Plattform (Heimvorteil)",
    "20% → Gemeinde",
    "10% → Dorffonds",
], 18)
add_text_box(slide, 7.5, 2.0, 5, 0.5, "Unit Economics", 22, True, DARK)
add_bullet_list(slide, 7.5, 2.7, 5, 4, [
    "Break-even: ~45 Abonnenten pro Dorf",
    "Ziel: 300 Abonnenten = ~€30k/Jahr pro Dorf",
    "Gemeinde-Anteil: ~€8.600/Jahr",
    "Dorffonds: ~€4.300/Jahr",
], 18)

# ========== SLIDE 6: How it Works ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, EMERALD_DARK)
add_text_box(slide, 1.5, 0.8, 10, 0.8, "So funktioniert's", 40, True, WHITE)
add_text_box(slide, 1.5, 2.5, 3, 0.5, "1. Abo wählen", 28, True, RGBColor(0x6E, 0xE7, 0xB7))
add_text_box(slide, 1.5, 3.2, 3, 1, "Bewohner wählen online ihr Heimvorteil-Abo ab €9,90/Monat", 16, False, RGBColor(0xA7, 0xF3, 0xD0))
add_text_box(slide, 5.5, 2.5, 3, 0.5, "2. Karte zeigen", 28, True, RGBColor(0x6E, 0xE7, 0xB7))
add_text_box(slide, 5.5, 3.2, 3, 1, "QR-Code am Handy oder NFC-Karte beim Einkaufen vorzeigen", 16, False, RGBColor(0xA7, 0xF3, 0xD0))
add_text_box(slide, 9.5, 2.5, 3, 0.5, "3. Sofort sparen", 28, True, RGBColor(0x6E, 0xE7, 0xB7))
add_text_box(slide, 9.5, 3.2, 3, 1, "15–25% Rabatt wird direkt abgezogen. Bezahlung wie gewohnt.", 16, False, RGBColor(0xA7, 0xF3, 0xD0))
add_text_box(slide, 1.5, 5.0, 10, 1, "Checkout dauert 15 Sekunden. Kein neues Kassensystem nötig.", 20, False, WHITE, PP_ALIGN.CENTER)

# ========== SLIDE 7: Traction ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1.5, 0.8, 10, 0.8, "Traction", 40, True, EMERALD_DARK)
add_bullet_list(slide, 1.5, 2.0, 10, 5, [
    "MVP gebaut und live (Next.js + PostgreSQL + Mollie)",
    "Pilot-Dorf: [Name] — [X] Betriebe angemeldet",
    "[X] Abonnenten in den ersten [Y] Wochen",
    "Key Metrics: Transaktionen/Woche, Retention, NPS",
    "Tech-Stack: Skalierbar auf 500+ Dörfer ohne Code-Änderungen",
], 20)
add_text_box(slide, 1.5, 5.5, 10, 1, "→ Platzhalter — wird mit echten Daten nach Pilotstart gefüllt", 14, True, GRAY, PP_ALIGN.LEFT)

# ========== SLIDE 8: Unit Economics ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1.5, 0.8, 10, 0.8, "Unit Economics", 40, True, EMERALD_DARK)
add_text_box(slide, 1.5, 2.0, 5, 0.5, "Kosten", 24, True, DARK)
add_bullet_list(slide, 1.5, 2.7, 5, 2.5, [
    "CAC: ~€5–10 (Mundpropaganda + Gemeinde-Newsletter)",
    "Hosting: €8,50/Mo (Hetzner VPS)",
    "Mollie Gebühren: ~€0,35 + 1,8% pro Transaktion",
], 17)
add_text_box(slide, 7.5, 2.0, 5, 0.5, "Erträge", 24, True, DARK)
add_bullet_list(slide, 7.5, 2.7, 5, 2.5, [
    "LTV: €150+ (Ø 12+ Monate × €12,90/Mo)",
    "LTV/CAC Ratio: 15–30x",
    "Bruttomarge: ~85%",
], 17)
add_text_box(slide, 1.5, 5.0, 10, 0.5, "Skalierung", 24, True, DARK)
add_bullet_list(slide, 1.5, 5.7, 10, 1.5, [
    "Grenzkosten pro Dorf: €5.000 (Pilot) → €200 (Dorf 50+)",
    "Template-Modell: gleicher Code, neues Dorf = neuer Tenant",
], 17)

# ========== SLIDE 9: Competitive Advantage ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1.5, 0.8, 10, 0.8, "Wettbewerbsvorteil", 40, True, EMERALD_DARK)
add_bullet_list(slide, 1.5, 2.0, 10, 5, [
    "Gemeinde-Partnerschaft = Distribution Moat (kein Wettbewerber bekommt Gemeinde-Newsletter)",
    "Keine Zahlungslizenz nötig (identifikationsbasiert, nicht zahlungsbasiert)",
    "Genossenschafts-DNA der DACH-Region = kulturelle Passung",
    "Daten-Insights für Gemeinden = politisches Buy-in",
    "Dorffonds = emotionales Argument für Bürgermeister und Bewohner",
], 20)
add_text_box(slide, 1.5, 5.5, 10, 1.5,
    "Heimvorteil  vs  Dorftaler (3-5% Bonus, Papier)  vs  Payback (0,5% Punkte)  vs  Groupon (einmalig, 50%+ Marge)",
    16, False, GRAY, PP_ALIGN.CENTER)

# ========== SLIDE 10: Team ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1.5, 0.8, 10, 0.8, "Team", 40, True, EMERALD_DARK)
add_text_box(slide, 1.5, 2.5, 10, 1, "[Gründer/in]", 28, True, DARK, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.5, 10, 1, "Technischer Founder — baut und liefert", 20, False, GRAY, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.0, 10, 0.5, "Beirat (geplant)", 22, True, DARK, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.6, 10, 1, "Bürgermeister · Handwerkskammer · Raiffeisenbank", 18, False, GRAY, PP_ALIGN.CENTER)

# ========== SLIDE 11: Funding & Roadmap ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 1.5, 0.8, 10, 0.8, "Finanzierung & Roadmap", 40, True, EMERALD_DARK)
add_text_box(slide, 1.5, 2.0, 5, 0.5, "Finanzierung", 22, True, DARK)
add_bullet_list(slide, 1.5, 2.7, 5, 3, [
    "Bootstrap: ~€2.000 (Domain, GmbH, NFC-Karten)",
    "AWS PreSeed: bis €200k (80% gefördert)",
    "FFG Kleinprojekt: bis €88.500",
    "Startnext Crowdfunding: €15–30k",
    "Seed: €200–500k (Impact Investoren)",
], 16)
add_text_box(slide, 7.5, 2.0, 5, 0.5, "Roadmap", 22, True, DARK)
add_bullet_list(slide, 7.5, 2.7, 5, 3, [
    "Mo 1–3: 1 Dorf, 200 Abonnenten",
    "Mo 4–9: 5 Dörfer, 1.000 Abonnenten",
    "Mo 10–18: 50 Dörfer, 5.000 Abonnenten",
    "Jahr 2+: 100+ Dörfer, AT/DE/CH",
    "Profitabilität: 10 Dörfer × 300 Abos = €300k ARR",
], 16)

# ========== SLIDE 12: The Ask ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, EMERALD_DARK)
add_text_box(slide, 1.5, 1.0, 10, 0.8, "The Ask", 40, True, WHITE)
add_text_box(slide, 1.5, 2.5, 10, 0.8, "€[X]k für 12 Monate", 36, True, RGBColor(0x6E, 0xE7, 0xB7), PP_ALIGN.CENTER)
add_bullet_list(slide, 2.5, 3.8, 8, 2.5, [
    "60% Produkt (Features, Mobile App, POS-Integration)",
    "20% Village Acquisition (Community Manager, Onboarding)",
    "20% Operations (Hosting, Legal, Admin)",
], 20, RGBColor(0xA7, 0xF3, 0xD0))
add_text_box(slide, 1.5, 5.8, 10, 0.5, "Ziel: 5-Dorf-Modell in 12 Monaten beweisen", 22, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 6.5, 10, 0.5, "heimvorteil.at · [name]@heimvorteil.at", 16, False, RGBColor(0xA7, 0xF3, 0xD0), PP_ALIGN.CENTER)

# Save
output_dir = "/Users/martin_pletzenauer/villagelalternative/docs/pitch"
pptx_path = os.path.join(output_dir, "Heimvorteil-Investor-Pitch.pptx")
prs.save(pptx_path)
print(f"Saved: {pptx_path}")
